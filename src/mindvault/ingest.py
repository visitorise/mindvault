"""External source ingestion — files, URLs, directories."""

from __future__ import annotations

import json
import re
import shutil
import subprocess
import sys
import urllib.request
import urllib.error
from pathlib import Path

from mindvault.llm import detect_llm, call_llm, estimate_cost, confirm_api_usage

_EXTRACTION_PROMPT = """Extract key concepts and relationships from this text.
Return JSON only:
{
  "nodes": [{"id": "slug_name", "label": "Human Name", "file_type": "document", "source_file": "path"}],
  "edges": [{"source": "id1", "target": "id2", "relation": "references|implements|related_to", "confidence": "EXTRACTED|INFERRED", "confidence_score": 0.8}]
}

Rules:
- Extract named concepts, entities, technologies, decisions
- EXTRACTED: explicitly stated relationship
- INFERRED: reasonable inference
- Keep nodes under 30 per document
- Keep edges under 50 per document"""


def _extract_text_from_file(file_path: Path) -> str | None:
    """Extract text content from a file based on extension."""
    ext = file_path.suffix.lower()

    if ext in (".md", ".txt", ".rst"):
        try:
            return file_path.read_text(encoding="utf-8", errors="ignore")
        except (OSError, IOError):
            return None

    if ext == ".pdf":
        return _extract_pdf_text(file_path)

    if ext == ".docx":
        return _extract_docx_text(file_path)

    if ext == ".xlsx":
        return _extract_xlsx_text(file_path)

    if ext == ".pptx":
        return _extract_pptx_text(file_path)

    if ext in (".png", ".jpg", ".jpeg", ".webp", ".gif"):
        # Image: skip (vision API is a Known Gap)
        return None

    # Unknown extension: try reading as text
    try:
        return file_path.read_text(encoding="utf-8", errors="ignore")
    except (OSError, IOError, UnicodeDecodeError):
        return None


def _extract_docx_text(file_path: Path) -> str | None:
    """Extract text from .docx. Requires `python-docx` (install via [office] extra)."""
    try:
        from docx import Document
    except ImportError:
        return None
    try:
        doc = Document(str(file_path))
        parts = [p.text for p in doc.paragraphs if p.text]
        for table in doc.tables:
            for row in table.rows:
                parts.append("\t".join(cell.text for cell in row.cells))
        return "\n".join(parts) if parts else None
    except Exception:
        return None


def _extract_xlsx_text(file_path: Path) -> str | None:
    """Extract text from .xlsx. Requires `openpyxl` (install via [office] extra)."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return None
    try:
        wb = load_workbook(str(file_path), data_only=True, read_only=True)
        lines = []
        for ws in wb.worksheets:
            lines.append(f"# Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = ["" if c is None else str(c) for c in row]
                if any(cells):
                    lines.append("\t".join(cells))
        wb.close()
        return "\n".join(lines) if lines else None
    except Exception:
        return None


def _extract_pptx_text(file_path: Path) -> str | None:
    """Extract text from .pptx. Requires `python-pptx` (install via [office] extra)."""
    try:
        from pptx import Presentation
    except ImportError:
        return None
    try:
        prs = Presentation(str(file_path))
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"# Slide {i}")
            for shape in slide.shapes:
                text = getattr(shape, "text", None)
                if text:
                    lines.append(text)
        return "\n".join(lines) if lines else None
    except Exception:
        return None


def _extract_pdf_text(file_path: Path) -> str | None:
    """Extract text from PDF using pdftotext, or skip."""
    try:
        result = subprocess.run(
            ["pdftotext", str(file_path), "-"],
            capture_output=True, text=True, timeout=30,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout
    except (FileNotFoundError, subprocess.TimeoutExpired, OSError):
        pass
    return None


def _strip_html(html: str) -> str:
    """Strip HTML to plain text (simple tag removal)."""
    # Remove script and style blocks
    html = re.sub(r"<script[^>]*>.*?</script>", "", html, flags=re.DOTALL | re.IGNORECASE)
    html = re.sub(r"<style[^>]*>.*?</style>", "", html, flags=re.DOTALL | re.IGNORECASE)
    # Remove tags
    html = re.sub(r"<[^>]+>", " ", html)
    # Collapse whitespace
    html = re.sub(r"\s+", " ", html).strip()
    return html


def _url_to_slug(url: str) -> str:
    """Convert URL to a filesystem-safe slug."""
    # Remove protocol
    slug = re.sub(r"^https?://", "", url)
    # Replace non-alphanumeric with underscore
    slug = re.sub(r"[^a-zA-Z0-9]", "_", slug)
    # Collapse underscores and trim
    slug = re.sub(r"_+", "_", slug).strip("_")
    return slug[:100]


def _llm_extract(text: str, source_file: str, provider: dict) -> dict:
    """Call LLM for concept extraction, return {nodes, edges} or empty."""
    if not text or not text.strip():
        return {"nodes": [], "edges": []}

    # Truncate text to max_tokens_per_file equivalent chars
    from mindvault.config import get as cfg_get
    max_tokens = cfg_get("max_tokens_per_file", 4000)
    max_chars = max_tokens * 4  # rough token-to-char ratio
    if len(text) > max_chars:
        text = text[:max_chars]

    response = call_llm(_EXTRACTION_PROMPT, text, provider)
    if not response:
        return {"nodes": [], "edges": []}

    # Parse JSON from response (may be wrapped in markdown code block)
    return _parse_llm_json(response, source_file)


def _parse_llm_json(response: str, source_file: str) -> dict:
    """Parse LLM JSON response, handling markdown code blocks.

    Node IDs emitted by the LLM are rewritten to canonical form
    ``{rel_path_slug}::concept::{slug}`` so ingest-generated nodes merge
    cleanly with extract_semantic output.
    """
    from mindvault.extract import _make_canonical_id, _make_ref_id

    # Strip markdown code block if present
    cleaned = response.strip()
    if cleaned.startswith("```"):
        # Remove first and last lines (```json and ```)
        lines = cleaned.split("\n")
        if lines[0].startswith("```"):
            lines = lines[1:]
        if lines and lines[-1].strip() == "```":
            lines = lines[:-1]
        cleaned = "\n".join(lines)

    try:
        data = json.loads(cleaned)
        nodes = data.get("nodes", [])
        edges = data.get("edges", [])

        # Rewrite LLM-chosen IDs into canonical form, then rewire edges.
        id_map: dict[str, str] = {}
        for node in nodes:
            old_id = node.get("id", "")
            label = node.get("label", old_id)
            new_id = _make_canonical_id(source_file, "concept", label or old_id)
            if old_id:
                id_map[old_id] = new_id
            node["id"] = new_id
            if "source_file" not in node or not node["source_file"]:
                node["source_file"] = source_file
            if "file_type" not in node:
                node["file_type"] = "document"
            node.setdefault("entity_type", "concept")

        # Ensure required edge fields + rewire IDs
        for edge in edges:
            src = edge.get("source", "")
            tgt = edge.get("target", "")
            edge["source"] = id_map.get(src, _make_ref_id(src) if src else src)
            edge["target"] = id_map.get(tgt, _make_ref_id(tgt) if tgt else tgt)
            if "confidence" not in edge:
                edge["confidence"] = "INFERRED"
            if "confidence_score" not in edge:
                edge["confidence_score"] = 0.7
            if "source_file" not in edge:
                edge["source_file"] = source_file
            if "weight" not in edge:
                edge["weight"] = 1.0

        return {"nodes": nodes, "edges": edges}
    except (json.JSONDecodeError, KeyError, TypeError):
        return {"nodes": [], "edges": []}


def _classify_into_communities(nodes: list[dict], concepts_path: Path, wiki_dir: Path) -> dict:
    """Classify extracted nodes into existing communities or new pages.

    Matches ingested node labels against existing _concepts.json entries.
    If match found, the node merges into the existing wiki page.
    If no match, the node goes to wiki/ingested/ as a new page.

    Args:
        nodes: List of extracted node dicts with 'label' key.
        concepts_path: Path to _concepts.json.
        wiki_dir: Wiki directory path.

    Returns:
        Dict with 'merged' (list of {node, target, concept}) and 'new' (list of nodes).
    """
    concepts: dict[str, list[str]] = {}
    if concepts_path.exists():
        try:
            concepts = json.loads(concepts_path.read_text(encoding="utf-8"))
        except (json.JSONDecodeError, OSError):
            concepts = {}

    classified: dict[str, list] = {"merged": [], "new": []}

    for node in nodes:
        label = node.get("label", "").lower()
        if not label:
            classified["new"].append(node)
            continue

        # Find best matching concept by word overlap
        best_match = None
        best_score = 0
        label_words = set(label.split())
        for concept, pages in concepts.items():
            concept_words = set(concept.split())
            overlap = len(label_words & concept_words)
            if overlap > best_score:
                best_score = overlap
                best_match = (concept, pages)

        if best_match and best_score > 0:
            target_page = best_match[1][0]  # First related page
            classified["merged"].append({
                "node": node,
                "target": target_page,
                "concept": best_match[0],
            })
        else:
            classified["new"].append(node)

    return classified


def _update_wiki_from_extraction(extraction: dict, file_path: Path, output_dir: Path) -> dict:
    """Update wiki pages based on LLM extraction results.

    Returns:
        Dict with 'merged_to_existing' and 'new_pages' counts.
    """
    nodes = extraction.get("nodes", [])
    if not nodes:
        return {"merged_to_existing": 0, "new_pages": 0}

    wiki_dir = output_dir / "wiki"
    ingested_dir = wiki_dir / "ingested"
    ingested_dir.mkdir(parents=True, exist_ok=True)

    # Load existing concepts index
    concepts_path = wiki_dir / "_concepts.json"
    concepts: dict[str, list[str]] = {}
    if concepts_path.exists():
        try:
            concepts = json.loads(concepts_path.read_text(encoding="utf-8"))
        except (json.JSONDecodeError, OSError):
            concepts = {}

    source_name = Path(file_path).stem

    # Classify nodes into existing communities or new pages
    classified = _classify_into_communities(nodes, concepts_path, wiki_dir)
    merged_count = 0
    new_count = 0

    # Process merged nodes — append to existing wiki pages
    for entry in classified["merged"]:
        node = entry["node"]
        target_page = entry["target"]
        page_path = wiki_dir / target_page
        if not page_path.exists():
            # Target page gone — treat as new
            classified["new"].append(node)
            continue
        content = page_path.read_text(encoding="utf-8")
        source_line = f"- {source_name} ({file_path})"
        if source_line in content:
            merged_count += 1
            continue  # already there
        if "## Ingested Sources" in content:
            content = content.replace(
                "## Ingested Sources",
                f"## Ingested Sources\n{source_line}",
            )
        else:
            marker = "<!-- user-notes -->"
            if marker in content:
                content = content.replace(
                    marker,
                    f"\n## Ingested Sources\n{source_line}\n\n{marker}",
                )
            else:
                content += f"\n## Ingested Sources\n{source_line}\n"
        page_path.write_text(content, encoding="utf-8")
        merged_count += 1

    # Process new nodes — create pages in ingested/
    for node in classified["new"]:
        label = node.get("label", node.get("id", "unknown"))
        concept_key = label.lower()
        node_id = node.get("id", label)

        slug = re.sub(r"[^a-z0-9\s-]", "", label.lower().strip())
        slug = re.sub(r"[\s]+", "-", slug)
        slug = re.sub(r"-+", "-", slug).strip("-")
        if not slug:
            slug = node_id

        page_path = ingested_dir / f"{slug}.md"
        page_content = [
            f"# {label}",
            f"Source: {file_path}",
            "",
            "## Ingested Sources",
            f"- {source_name} ({file_path})",
            "",
        ]

        # Add edges info if available
        edges = extraction.get("edges", [])
        related = [e for e in edges if e.get("source") == node_id or e.get("target") == node_id]
        if related:
            page_content.append("## Relationships")
            for edge in related:
                src = edge.get("source", "")
                tgt = edge.get("target", "")
                rel = edge.get("relation", "related")
                if src == node_id:
                    page_content.append(f"- -> {rel} -> {tgt}")
                else:
                    page_content.append(f"- <- {rel} <- {src}")

        page_path.write_text("\n".join(page_content) + "\n", encoding="utf-8")

        # Register in concepts
        rel_page = f"ingested/{slug}.md"
        if concept_key not in concepts:
            concepts[concept_key] = []
        if rel_page not in concepts[concept_key]:
            concepts[concept_key].append(rel_page)
        new_count += 1

    # Write updated concepts index
    concepts_path.parent.mkdir(parents=True, exist_ok=True)
    concepts_path.write_text(json.dumps(concepts, ensure_ascii=False, indent=2), encoding="utf-8")

    # Create ingested/INDEX.md
    ingested_pages = sorted(ingested_dir.glob("*.md"))
    index_lines = ["# Ingested Sources", ""]
    for p in ingested_pages:
        if p.name == "INDEX.md":
            continue
        index_lines.append(f"- [[{p.stem}]]")
    (ingested_dir / "INDEX.md").write_text("\n".join(index_lines) + "\n", encoding="utf-8")

    # Update search index
    _update_search_index_for_ingested(wiki_dir, output_dir)

    return {"merged_to_existing": merged_count, "new_pages": new_count}


def _update_search_index_for_ingested(wiki_dir: Path, output_dir: Path) -> None:
    """Update search index with ingested wiki pages."""
    index_path = output_dir / "search_index.json"
    if not index_path.exists():
        return

    from mindvault.index import load_index, _tokenize, _extract_title, _extract_headings, _hash_content, _compute_idf

    index_data = load_index(index_path)
    docs = index_data.get("docs", {})

    ingested_dir = wiki_dir / "ingested"
    if not ingested_dir.exists():
        return

    for md_file in ingested_dir.glob("*.md"):
        content = md_file.read_text(encoding="utf-8")
        rel_path = str(md_file.relative_to(wiki_dir))
        docs[rel_path] = {
            "title": _extract_title(content) or md_file.stem,
            "headings": _extract_headings(content),
            "tokens": _tokenize(content),
            "hash": _hash_content(content),
        }

    index_data["docs"] = docs
    index_data["doc_count"] = len(docs)
    index_data["idf"] = _compute_idf(docs)
    index_path.write_text(json.dumps(index_data, ensure_ascii=False, indent=2), encoding="utf-8")


def ingest_file(file_path: Path, output_dir: Path) -> dict:
    """Ingest a single file: copy to sources/, extract text, LLM extract, merge.

    Returns:
        Dict with keys: nodes (int), edges (int), source (str), or {skipped: True}.
    """
    file_path = Path(file_path).resolve()
    output_dir = Path(output_dir)

    if not file_path.exists():
        return {"error": f"File not found: {file_path}"}

    # Copy to sources/
    sources_dir = output_dir / "sources"
    sources_dir.mkdir(parents=True, exist_ok=True)
    dest = sources_dir / file_path.name
    shutil.copy2(file_path, dest)

    # Extract text
    text = _extract_text_from_file(file_path)
    if text is None:
        return {"skipped": True, "source": str(file_path), "reason": "unsupported format"}

    # Detect LLM
    provider = detect_llm()
    if provider["provider"] is None:
        # No LLM available — still create wiki page from file metadata.
        # Use canonical file-level node ID so graph-side merges stay stable.
        from mindvault.extract import _make_canonical_id
        fallback_id = _make_canonical_id(str(file_path), "file", file_path.stem)
        wiki_result = _update_wiki_from_extraction(
            {
                "nodes": [{
                    "id": fallback_id,
                    "label": file_path.stem,
                    "entity_type": "file",
                    "source_file": str(file_path),
                }],
                "edges": [],
            },
            file_path, output_dir,
        )
        return {
            "nodes": 1, "edges": 0, "source": str(file_path),
            "reason": "no LLM, metadata only",
            "merged_to_existing": wiki_result.get("merged_to_existing", 0),
            "new_pages": wiki_result.get("new_pages", 0),
        }

    # Consent for API
    if not provider["is_local"]:
        cost = estimate_cost(text, provider)
        if not confirm_api_usage(provider, cost):
            return {"skipped": True, "source": str(file_path), "reason": "API usage declined"}

    # LLM extraction
    result = _llm_extract(text, str(file_path), provider)

    # Update wiki with extracted concepts
    wiki_result = _update_wiki_from_extraction(result, file_path, output_dir)

    return {
        "nodes": len(result["nodes"]),
        "edges": len(result["edges"]),
        "source": str(file_path),
        "extraction": result,
        "merged_to_existing": wiki_result.get("merged_to_existing", 0),
        "new_pages": wiki_result.get("new_pages", 0),
    }


def ingest_url(url: str, output_dir: Path) -> dict:
    """Ingest a URL: fetch, strip HTML, save as .md, then extract.

    Returns:
        Dict with keys: nodes (int), edges (int), source (str).
    """
    output_dir = Path(output_dir)
    sources_dir = output_dir / "sources"
    sources_dir.mkdir(parents=True, exist_ok=True)

    # Fetch URL
    try:
        req = urllib.request.Request(
            url,
            headers={"User-Agent": "MindVault/0.1 (knowledge extraction)"},
        )
        resp = urllib.request.urlopen(req, timeout=30)
        html = resp.read().decode("utf-8", errors="ignore")
    except Exception as e:
        return {"error": f"Failed to fetch URL: {e}"}

    # Strip HTML to text
    text = _strip_html(html)

    # Save as markdown
    slug = _url_to_slug(url)
    md_path = sources_dir / f"{slug}.md"
    md_path.write_text(f"# {url}\n\n{text}", encoding="utf-8")

    # Now ingest the saved file
    return ingest_file(md_path, output_dir)


def ingest(source: str, output_dir: Path) -> dict:
    """Entry point: detect file vs URL vs directory and dispatch.

    Args:
        source: File path, URL, or directory path.
        output_dir: MindVault output directory.

    Returns:
        Dict with ingestion results.
    """
    output_dir = Path(output_dir)

    # URL detection
    if source.startswith("http://") or source.startswith("https://"):
        return ingest_url(source, output_dir)

    path = Path(source)
    if not path.exists():
        return {"error": f"Source not found: {source}"}

    if path.is_file():
        return ingest_file(path, output_dir)

    if path.is_dir():
        # Recursively ingest all supported files, honoring SKIP_DIRS (incl. .obsidian)
        from mindvault.detect import SKIP_DIRS
        import os as _os

        total_nodes = 0
        total_edges = 0
        files_processed = 0

        for dirpath, dirnames, filenames in _os.walk(path):
            # Prune traversal in-place — skip VCS, build artifacts, Obsidian internals, etc.
            dirnames[:] = sorted(d for d in dirnames if d not in SKIP_DIRS and not d.startswith("."))
            for fname in sorted(filenames):
                if fname.startswith("."):
                    continue
                child = Path(dirpath) / fname
                result = ingest_file(child, output_dir)
                if not result.get("skipped") and not result.get("error"):
                    total_nodes += result.get("nodes", 0)
                    total_edges += result.get("edges", 0)
                    files_processed += 1
        return {
            "nodes": total_nodes,
            "edges": total_edges,
            "files_processed": files_processed,
            "source": str(path),
        }

    return {"error": f"Unknown source type: {source}"}
